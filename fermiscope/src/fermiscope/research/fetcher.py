"""DocumentFetcher — 安全なURL取得と文書パース。

- 取得前後のURL安全性検査(リダイレクト先も再検査)
- 応答サイズ上限(ストリーミングで検査)
- Content-Type 許可リスト
- robots.txt の尊重(取得失敗時は許可として続行)
- HTML/PDF/CSV/JSON のテキスト・表抽出(JSは実行しない)
- インメモリキャッシュ(TTL付き)
"""

from __future__ import annotations

import hashlib
import io
import ipaddress
import re
import time
import urllib.robotparser
from dataclasses import dataclass, field
from datetime import UTC, datetime
from urllib.parse import urljoin, urlparse

import httpx

from fermiscope.config import Settings
from fermiscope.domain.enums import DocumentType
from fermiscope.security.sanitizer import sanitize_extracted_text, strip_html_to_text
from fermiscope.security.url_guard import (
    Resolver,
    UrlGuardError,
    default_resolver,
    validate_ip,
    validate_url,
)


class FetchError(RuntimeError):
    pass


@dataclass
class FetchedDocument:
    url: str
    final_url: str
    content_type: str
    doc_type: DocumentType
    status_code: int
    text: str = ""
    html: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)
    title: str = ""
    content_hash: str = ""
    fetched_at: datetime = field(default_factory=lambda: datetime.now(UTC))
    size_bytes: int = 0


_OOXML_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocumentType.DOCX,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": DocumentType.XLSX,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": DocumentType.PPTX,
}
_OOXML_EXTENSIONS = {
    ".docx": DocumentType.DOCX,
    ".xlsx": DocumentType.XLSX,
    ".pptx": DocumentType.PPTX,
}
# 実サーバが office/pdf を配る際に使いがちな汎用バイナリ Content-Type。
# これらのときだけ拡張子ベースの判定を許す(空ヘッダは含めない)。
_GENERIC_BINARY_TYPES = {
    "application/octet-stream",
    "application/zip",
    "application/x-zip-compressed",
    "binary/octet-stream",
}


def _detect_doc_type(content_type: str, url: str) -> DocumentType:
    ct = content_type.split(";")[0].strip().lower()
    path = urlparse(url).path.lower()
    if ct in ("text/html", "application/xhtml+xml"):
        return DocumentType.HTML
    if ct in _OOXML_CONTENT_TYPES:
        return _OOXML_CONTENT_TYPES[ct]
    for ext, dt in _OOXML_EXTENSIONS.items():
        if path.endswith(ext):
            return dt
    if ct == "application/pdf" or path.endswith(".pdf"):
        return DocumentType.PDF
    if ct == "text/csv" or path.endswith(".csv"):
        return DocumentType.CSV
    if ct == "application/json" or path.endswith(".json"):
        return DocumentType.JSON
    if ct.startswith("text/"):
        return DocumentType.TEXT
    return DocumentType.UNKNOWN


def _extract_html_tables(html: str) -> list[list[list[str]]]:
    """BeautifulSoupでHTML表をセル文字列の行列として抽出する。"""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    tables: list[list[list[str]]] = []
    for table in soup.find_all("table"):
        rows: list[list[str]] = []
        for tr in table.find_all("tr"):
            cells = [c.get_text(" ", strip=True) for c in tr.find_all(["th", "td"])]
            if cells:
                rows.append(cells)
        if rows:
            tables.append(rows)
    return tables


def _extract_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: S112 — 壊れたページは飛ばして残りを抽出する
            continue
    return "\n".join(parts).strip()


def _guard_zip_bomb(data: bytes, max_uncompressed: int) -> None:
    """Office(ZIP)文書の解凍後合計サイズを検査し、ZIP爆弾を弾く。

    実際に解凍せず、ZIPディレクトリが申告する非圧縮サイズの合計で判定する。
    ZIPとして開けない場合はここでは何もしない(後段のパーサが弾く)。
    """
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except (zipfile.BadZipFile, OSError):
        return
    if total > max_uncompressed:
        raise FetchError(
            f"文書の解凍後サイズが大きすぎます({total} > {max_uncompressed} バイト)"
        )


def _extract_docx_text(data: bytes) -> str:
    """.docx(Word)から本文・表テキストを抽出する(マクロ・外部参照は実行しない)。"""
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for para in document.paragraphs[:50000]:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables[:500]:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" ".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(data: bytes) -> tuple[str, list[list[list[str]]]]:
    """.xlsx(Excel)から表(セル値)とテキストを抽出する。

    read_only=True でストリーミング読み出し、data_only=True で数式は評価せず
    保存済みの計算結果値のみを読む(数式文字列の露出・評価を避ける)。
    """
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    tables: list[list[list[str]]] = []
    text_parts: list[str] = []
    try:
        for ws in wb.worksheets[:50]:
            rows: list[list[str]] = []
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if r_idx >= 5000:
                    break
                cells = ["" if v is None else str(v) for v in row[:200]]
                if any(c.strip() for c in cells):
                    rows.append(cells)
                    text_parts.append(" ".join(cells))
            if rows:
                tables.append(rows)
    finally:
        wb.close()
    return "\n".join(text_parts).strip(), tables


def _extract_pptx_text(data: bytes) -> str:
    """.pptx(PowerPoint)からスライド上の可視テキスト・表を抽出する。

    ノート(発表者メモ)は通常の閲覧者に見えず隠しテキストの温床になるため抽出しない。
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in list(prs.slides)[:1000]:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" ".join(cells))
    return "\n".join(parts).strip()


def _html_title(html: str) -> str:
    m = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
    return m.group(1).strip()[:300] if m else ""


class DocumentFetcher:
    def __init__(
        self,
        settings: Settings,
        transport: httpx.AsyncBaseTransport | None = None,
        resolver: Resolver | None = None,
        skip_dns: bool = False,
        respect_robots: bool = True,
    ) -> None:
        self.settings = settings
        self._transport = transport
        self._resolver = resolver
        self._skip_dns = skip_dns
        self._respect_robots = respect_robots
        self._cache: dict[str, tuple[float, FetchedDocument]] = {}
        self._robots_cache: dict[str, urllib.robotparser.RobotFileParser | None] = {}
        self._client = httpx.AsyncClient(
            transport=transport,
            follow_redirects=False,  # リダイレクトは手動で辿り、各ホップを再検査する
            timeout=settings.fetch.timeout_seconds,
            headers={"User-Agent": settings.fetch.user_agent},
        )

    async def close(self) -> None:
        await self._client.aclose()

    def _pin_connection(self, url: str) -> tuple[str, str | None, str | None]:
        """ホスト名を1回だけ解決して検証済みIPへ接続を固定する(DNSリバインディング対策)。

        検証時と接続時でDNSが差し替わる TOCTOU を防ぐため、解決したIPで接続URLを
        書き換え、Host ヘッダと TLS SNI は元のホスト名に保つ。

        Returns:
            (接続URL, Hostヘッダ, SNIホスト名)。書き換え不要なら (url, None, None)。
        """
        # モック(オフライン)モードでは書き換えない — フィクスチャは元URLで配信される
        if self._skip_dns:
            return url, None, None
        parsed = urlparse(url)
        host = parsed.hostname or ""
        # 既にIPリテラル(validate_url で検証済み)ならそのまま
        try:
            ipaddress.ip_address(host.strip("[]"))
            return url, None, None
        except ValueError:
            pass
        ips = (self._resolver or default_resolver)(host)
        if not ips:
            raise UrlGuardError(f"ホスト名を解決できません: {host}")
        for ip in ips:
            validate_ip(ip)  # 全解決IPを検査
        pinned = ips[0]
        ip_host = f"[{pinned}]" if ":" in pinned else pinned
        port = f":{parsed.port}" if parsed.port else ""
        connect_url = parsed._replace(netloc=f"{ip_host}{port}").geturl()
        return connect_url, parsed.netloc, host

    async def _pinned_send(self, url: str, *, stream: bool) -> httpx.Response:
        """検証済みIPへ固定して送信する。"""
        connect_url, host_header, sni = self._pin_connection(url)
        request = self._client.build_request("GET", connect_url)
        if host_header:
            request.headers["Host"] = host_header
        if sni:
            request.extensions["sni_hostname"] = sni  # TLS証明書は元ホスト名で検証
        return await self._client.send(request, stream=stream)

    async def _read_capped(self, resp: httpx.Response, max_bytes: int) -> bytes:
        """ストリーム応答を上限つきで読み切る。超過時は FetchError。"""
        declared = resp.headers.get("content-length")
        if declared:
            try:
                if int(declared) > max_bytes:
                    raise FetchError(f"応答サイズが上限({max_bytes}バイト)を超えています")
            except ValueError:
                pass  # 不正な Content-Length はヘッダ値を無視しストリームで実測検査する
        chunks: list[bytes] = []
        total = 0
        async for chunk in resp.aiter_bytes():
            total += len(chunk)
            if total > max_bytes:
                raise FetchError(f"応答サイズが上限({max_bytes}バイト)を超えました")
            chunks.append(chunk)
        return b"".join(chunks)

    async def _robots_allows(self, url: str) -> bool:
        if not self._respect_robots:
            return True
        parsed = urlparse(url)
        base = f"{parsed.scheme}://{parsed.netloc}"
        if base not in self._robots_cache:
            rp = urllib.robotparser.RobotFileParser()
            try:
                resp = await self._pinned_send(base + "/robots.txt", stream=True)
                try:
                    # robots.txt も応答サイズ上限を適用(巨大 robots によるDoS防止)
                    cap = min(self.settings.fetch.max_response_bytes, 1024 * 1024)
                    body = await self._read_capped(resp, cap)
                    status = resp.status_code
                finally:
                    await resp.aclose()
                if status == 200:
                    rp.parse(body.decode("utf-8", errors="replace").splitlines())
                    self._robots_cache[base] = rp
                else:
                    self._robots_cache[base] = None  # robotsなし → 許可
            except (httpx.HTTPError, UrlGuardError, FetchError):
                # robots取得先が解決不能/プライベートIP/巨大でも許可扱い(本体fetchで再検査)
                self._robots_cache[base] = None
        rp2 = self._robots_cache[base]
        if rp2 is None:
            return True
        return rp2.can_fetch(self.settings.fetch.user_agent, url)

    async def _get_limited(self, url: str) -> httpx.Response:
        """サイズ上限つきGET。超過時は打ち切って FetchError。接続は検証済みIPへ固定。"""
        max_bytes = self.settings.fetch.max_response_bytes
        resp = await self._pinned_send(url, stream=True)
        try:
            content = await self._read_capped(resp, max_bytes)
        finally:
            await resp.aclose()
        resp._content = content  # noqa: SLF001
        return resp

    async def fetch(self, url: str) -> FetchedDocument:
        """URLを安全に取得し、パース済み文書を返す。"""
        ttl = self.settings.fetch.cache_ttl_hours * 3600
        cached = self._cache.get(url)
        if cached and time.time() - cached[0] < ttl:
            return cached[1]

        current = validate_url(url, resolver=self._resolver, skip_dns=self._skip_dns)

        redirects = 0
        while True:
            # 各ホップで robots.txt を確認(リダイレクト先のポリシーも尊重する)
            if not await self._robots_allows(current):
                raise FetchError(f"robots.txt により取得が許可されていません: {current}")
            try:
                resp = await self._get_limited(current)
            except httpx.HTTPError as exc:
                raise FetchError(f"取得に失敗しました: {type(exc).__name__}") from exc
            if resp.status_code in (301, 302, 303, 307, 308):
                redirects += 1
                if redirects > self.settings.fetch.max_redirects:
                    raise FetchError("リダイレクト回数が上限を超えました")
                location = resp.headers.get("location", "")
                if not location:
                    raise FetchError("リダイレクト先がありません")
                next_url = urljoin(current, location)
                # リダイレクト先も再検査(SSRF対策)
                current = validate_url(next_url, resolver=self._resolver, skip_dns=self._skip_dns)
                continue
            break

        if resp.status_code != 200:
            raise FetchError(f"HTTP {resp.status_code}: {current}")

        content_type = resp.headers.get("content-type", "")
        base_ct = content_type.split(";")[0].strip().lower()
        allowed = self.settings.fetch.allowed_content_types
        data = resp.content
        doc_type = _detect_doc_type(content_type, current)
        # 実サーバは office/pdf を汎用バイナリ(octet-stream / zip)で配ることがある。
        # その場合に限り URL拡張子で文書型を判定し、実バイトをパーサに検証させる。
        # 空/未宣言の Content-Type は従来どおり拒否(拡張子だけでは通さない)。
        allow_by_extension = base_ct in _GENERIC_BINARY_TYPES and doc_type != DocumentType.UNKNOWN
        if base_ct not in allowed and not allow_by_extension:
            raise FetchError(
                f"許可されていない/未宣言のContent-Typeです: {base_ct or '(なし)'}"
            )

        doc = FetchedDocument(
            url=url,
            final_url=current,
            content_type=base_ct,
            doc_type=doc_type,
            status_code=resp.status_code,
            content_hash=hashlib.sha256(data).hexdigest()[:16],
            size_bytes=len(data),
        )
        max_chars = self.settings.fetch.max_extracted_chars
        if doc_type == DocumentType.HTML:
            html = data.decode(resp.encoding or "utf-8", errors="replace")
            doc.html = html
            doc.text = strip_html_to_text(html)  # scriptは除去され実行されない
            doc.tables = _extract_html_tables(html)
            doc.title = _html_title(html)
        elif doc_type == DocumentType.PDF:
            try:
                doc.text = _extract_pdf_text(data)
            except Exception as exc:
                raise FetchError(f"PDFのテキスト抽出に失敗しました: {type(exc).__name__}") from exc
        elif doc_type in (DocumentType.DOCX, DocumentType.XLSX, DocumentType.PPTX):
            _guard_zip_bomb(data, self.settings.fetch.max_office_uncompressed_bytes)
            try:
                if doc_type == DocumentType.DOCX:
                    doc.text = _extract_docx_text(data)
                elif doc_type == DocumentType.XLSX:
                    doc.text, doc.tables = _extract_xlsx(data)
                else:
                    doc.text = _extract_pptx_text(data)
            except FetchError:
                raise
            except Exception as exc:
                raise FetchError(
                    f"Office文書の抽出に失敗しました({doc_type.value}): {type(exc).__name__}"
                ) from exc
        elif doc_type in (DocumentType.CSV, DocumentType.JSON, DocumentType.TEXT):
            doc.text = data.decode(resp.encoding or "utf-8", errors="replace")
        else:
            raise FetchError(f"未対応の文書タイプです: {content_type}")

        # 外部文書由来の本文は不信データ。不可視・制御文字を除去し長さを制限する
        # (隠しプロンプトの無害化・多層防御)。タイトルも同様に無害化する。
        doc.text = sanitize_extracted_text(doc.text, max_chars)
        if doc.title:
            doc.title = sanitize_extracted_text(doc.title, 300)

        self._cache[url] = (time.time(), doc)
        return doc


__all__ = ["DocumentFetcher", "FetchError", "FetchedDocument", "UrlGuardError"]
