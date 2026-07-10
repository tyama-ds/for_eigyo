"""BookIndex — 論文 BookRAG (arXiv:2512.03413) の文書ネイティブ索引 B = (T, G, M)。

- T: 文書の論理階層を表す木（目次に相当）。ノードは Section / Text / Table / Image。
- G: ノードから抽出したエンティティと関係の知識グラフ（KG）。
- M: GT-Link。各エンティティを抽出元のツリーノード集合へ対応付ける。

本モジュールは BookIndex の「構築」を担う:
  4.2 Tree Construction（Layout Parsing → Section Filtering）
  4.3 Graph Construction（KG Construction → Gradient-based Entity Resolution / Algorithm 1）

論文との差分（環境前提による簡略化。本実装は OpenAI 互換API + JupyterLab 完結が前提）:
  - 版面解析に MinerU を使わず、Markdown/テキストの見出し or PDF テキストの
    ヒューリスティック + LLM Section Filtering で木を作る。
  - Rerank モデルの代わりに埋め込みコサイン類似度を Gradient-based ER のスコアに用いる
    （reranker=None 時）。画像は VLM ではなくキャプション/テキストとして扱う。
"""

from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field
from pathlib import Path

import numpy as np

# --------------------------------------------------------------------------
# LLM / 埋め込みの共通ヘルパー（llmlab の OpenAI 互換クライアントを使う）
# --------------------------------------------------------------------------


def llm_text(prompt: str, *, system: str | None = None, temperature: float = 0.0,
             max_retries: int | None = None, max_tokens: int | None = None) -> str:
    """max_retries=0 で SDK リトライを無効化（fail-fast）。上位で独自に再試行する用途向け。
    max_tokens で生成量を制限（遅いローカルLLMでの長時間生成→タイムアウトを防ぐ）。"""
    from .client import get_client
    from .config import get_settings

    s = get_settings()
    client = get_client()
    if max_retries is not None:
        try:
            client = client.with_options(max_retries=max_retries)
        except Exception:  # noqa: BLE001  古いSDK/スタブは with_options 非対応でもよい
            pass
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})
    kwargs = {"model": s.model, "messages": messages, "temperature": temperature}
    if max_tokens is not None:
        kwargs["max_tokens"] = max_tokens
    resp = client.chat.completions.create(**kwargs)
    raw = (resp.choices[0].message.content or "") if resp.choices else ""
    from .client import strip_think

    return strip_think(raw)  # 推論モデルの思考過程を除去


def llm_json(prompt: str, *, system: str | None = None, max_retries: int | None = None,
             max_tokens: int | None = None):
    """LLM 応答を JSON として解釈する（思考過程・```フェンス・前後テキストに耐える）。"""
    raw = llm_text(prompt, system=system, max_retries=max_retries, max_tokens=max_tokens)
    return parse_json_answer(raw)


def parse_json_answer(raw: str):
    """応答テキストから答えの JSON を取り出す。

    - 思考過程（<think>…</think> / 閉じタグのみ残る型）は llm_text 側で除去済みだが、
      タグ無しで思考を書き出すモデルにも耐えるよう、**最後に現れる完全なトップレベル
      JSON** を採用する（思考中に JSON 断片を書き散らしても、最終回答が勝つ）。
    - 旧実装の「最初の { から最後の } まで」の貪欲マッチは、思考混じりの応答で
      必ず parse 失敗していた。
    """
    if not raw:
        return None
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
        text = re.sub(r"\n?```$", "", text).strip()
    dec = json.JSONDecoder()
    last = None
    pos = 0
    while True:
        m = re.search(r"[\{\[]", text[pos:])
        if not m:
            break
        i = pos + m.start()
        try:
            obj, end = dec.raw_decode(text[i:])
            last = obj                 # 直近のトップレベル JSON を記憶
            pos = i + end              # その JSON の直後から続きを探す
        except json.JSONDecodeError:
            pos = i + 1
    return last


def vlm_describe(png_bytes: bytes, *, model: str | None = None) -> str:
    """画像を VLM（画像対応モデル）に渡して内容説明を得る。図が無ければ空文字。

    model 省略時は接続設定の model を使う（画像対応モデルであることが前提）。
    """
    import base64

    from .client import get_client
    from .config import get_settings

    s = get_settings()
    b64 = base64.b64encode(png_bytes).decode("ascii")
    resp = get_client().chat.completions.create(
        model=model or s.model,
        messages=[{"role": "user", "content": [
            {"type": "text", "text":
                "この画像の図・グラフ・表の内容と、読み取れる主要な数値・傾向を簡潔に"
                "説明してください。意味のある図が無ければ「図なし」とだけ答えてください。"},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
        ]}],
        max_tokens=400,
    )
    from .client import strip_think

    out = strip_think((resp.choices[0].message.content or "") if resp.choices else "")
    return "" if "図なし" in out else out.strip()


def embed(texts: list[str]) -> np.ndarray:
    """テキスト群を L2 正規化済みの埋め込み行列に変換する（embed_base_url を尊重）。"""
    from .client import get_embed_client
    from .config import get_settings

    if not texts:
        return np.zeros((0, 1), dtype=np.float32)
    s = get_settings()
    resp = get_embed_client().embeddings.create(model=s.embed_model, input=texts)
    vecs = np.array([d.embedding for d in resp.data], dtype=np.float32)
    norms = np.linalg.norm(vecs, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return vecs / norms


def cosine(a: np.ndarray, b: np.ndarray) -> float:
    na, nb = np.linalg.norm(a), np.linalg.norm(b)
    if na == 0 or nb == 0:
        return 0.0
    return float(np.dot(a, b) / (na * nb))


def progress(iterable, *, total: int | None = None, desc: str = ""):
    """tqdm があればプログレスバー、無ければ定期 print で進捗を返すイテレータ。

    JupyterLab では tqdm.auto がセル内にバーを描画する（現在のフェーズ=desc 付き）。
    """
    try:
        from tqdm.auto import tqdm

        return tqdm(iterable, total=total, desc=desc, leave=True)
    except Exception:  # noqa: BLE001
        def _gen():
            n = 0
            for x in iterable:
                n += 1
                if total and (n % 20 == 0 or n == total):
                    print(f"[BookRAG] {desc} {n}/{total}")
                elif not total and n % 20 == 0:
                    print(f"[BookRAG] {desc} {n}")
                yield x
        return _gen()


# ログ転送はスレッドローカル。以前は呼び出し側が bx.log をグローバルに差し替えて
# いたが、Studio はタスクごとに別スレッドで動くため、並行実行で転送先が混線し
# 復元順によっては死んだタスクの Queue を永久に指す事故が起きる。log_to() を使う。
_log_local = __import__("threading").local()


def log_to(callback):
    """このスレッドの log() 出力を callback にも転送するコンテキストマネージャ。

    使い方::
        with bx.log_to(lambda msg: progress_queue.put(msg)):
            book.add_book(...)

    ネスト可（内側が優先、抜けると外側に戻る）。他スレッドには影響しない。
    注意: ThreadPoolExecutor 内の子スレッドからの log() には転送されない
    （現状 log() を呼ぶのはタスクの主スレッドのみ）。
    """
    from contextlib import contextmanager

    @contextmanager
    def _cm():
        prev = getattr(_log_local, "cb", None)
        _log_local.cb = callback
        try:
            yield
        finally:
            _log_local.cb = prev

    return _cm()


def log(msg: str) -> None:
    cb = getattr(_log_local, "cb", None)
    if cb is not None:
        try:
            cb(msg)
        except Exception:  # noqa: BLE001  転送失敗で本処理は止めない
            pass
    print(f"[BookRAG] {msg}")


# --------------------------------------------------------------------------
# データ構造
# --------------------------------------------------------------------------

NODE_TYPES = ("Section", "Text", "Table", "Image")


@dataclass
class TreeNode:
    """ツリー T のノード。content と最終 type（4.2.2）を保持する。"""

    id: int
    type: str  # Section / Text / Table / Image
    content: str
    book: str
    level: int | None = None  # Section のみ意味を持つ（l=1 が最上位）
    title: str | None = None
    page: int | None = None
    parent: int | None = None
    children: list[int] = field(default_factory=list)
    # 元ファイルの絶対パス（書のルートノードにのみ設定。出典リンク用）。
    # 旧バージョンで保存した索引には無いため既定 None（load 互換）。
    source: str | None = None
    # 文書ID（内容ハッシュ）。title/book 名ではなくこれで文書を識別する。
    # 旧バージョンで保存した索引には無いため既定 None（load 互換）。
    doc_id: str | None = None


@dataclass
class Entity:
    """KG G の頂点。origin_nodes が GT-Link M（V→P(N)）を構成する。"""

    id: int
    name: str
    type: str
    description: str
    origin_nodes: list[int] = field(default_factory=list)


class BookIndex:
    """B = (T, G, M)。木・KG・GT-Link を一体で保持し、永続化できる。"""

    def __init__(self):
        self.nodes: dict[int, TreeNode] = {}
        self.roots: list[int] = []
        self.entities: dict[int, Entity] = {}
        self.relations: list[tuple[int, int, str]] = []  # (src_eid, dst_eid, label)
        self._node_seq = 0
        self._ent_seq = 0
        # エンティティ・ベクタDB（容量ダブリングの成長バッファ。append 償却 O(1)）
        self._entity_vecs: np.ndarray = np.zeros((0, 0), dtype=np.float32)
        self._entity_ids: list[int] = []  # 行 -> entity id
        self._row_of: dict[int, int] = {}  # entity id -> 行
        self._count = 0
        self._dim: int | None = None

    # ---- ノード ----
    def add_node(self, **kwargs) -> TreeNode:
        nid = self._node_seq
        self._node_seq += 1
        node = TreeNode(id=nid, **kwargs)
        self.nodes[nid] = node
        return node

    def node_to_entities(self) -> dict[int, list[int]]:
        """逆引き: ツリーノード → そこに紐づくエンティティ集合。"""
        rev: dict[int, list[int]] = {}
        for e in self.entities.values():
            for n in e.origin_nodes:
                rev.setdefault(n, []).append(e.id)
        return rev

    def subtree(self, node_id: int) -> list[int]:
        """node_id を根とする部分木の全ノード id（式(5) Subtree）。"""
        out, stack = [], [node_id]
        while stack:
            cur = stack.pop()
            out.append(cur)
            stack.extend(self.nodes[cur].children)
        return out

    # ---- エンティティ（KG） ----
    def _validate_dim(self, vec: np.ndarray) -> None:
        if self._dim is not None and vec.shape[0] != self._dim:
            raise RuntimeError(
                f"埋め込み次元が既存インデックス（{self._dim}）と一致しません（{vec.shape[0]}）。"
                "埋め込みモデルを変更した場合は、storage を削除（reset()）して作り直してください。"
            )

    def add_entity(self, name: str, etype: str, desc: str, origin: int, vec: np.ndarray) -> Entity:
        vec = _normalize(vec).astype(np.float32)
        self._validate_dim(vec)  # 変異前に検証（失敗時に entities とベクタが食い違わないように）
        eid = self._ent_seq
        self._ent_seq += 1
        ent = Entity(id=eid, name=name, type=etype, description=desc, origin_nodes=[origin])
        self.entities[eid] = ent
        self._append_vec(eid, vec)
        return ent

    def merge_entity(self, target_eid: int, name: str, origin: int, vec: np.ndarray) -> None:
        """新エンティティを既存 target にマージし、GT-Link を集約（4.3.2 末尾）。"""
        vec = _normalize(vec).astype(np.float32)
        self._validate_dim(vec)  # 変異前に検証
        ent = self.entities[target_eid]
        if origin not in ent.origin_nodes:
            ent.origin_nodes.append(origin)
        row = self._row_of[target_eid]  # O(1)
        self._entity_vecs[row] = _normalize((self._entity_vecs[row] + vec) / 2.0)

    def _append_vec(self, eid: int, vec: np.ndarray) -> None:
        vec = _normalize(vec).astype(np.float32)
        self._validate_dim(vec)
        if self._dim is None:
            self._dim = int(vec.shape[0])
            self._entity_vecs = np.zeros((64, self._dim), dtype=np.float32)
        if self._count >= self._entity_vecs.shape[0]:  # 容量ダブリング
            grown = np.zeros((self._entity_vecs.shape[0] * 2, self._dim), dtype=np.float32)
            grown[: self._count] = self._entity_vecs[: self._count]
            self._entity_vecs = grown
        self._entity_vecs[self._count] = vec
        self._row_of[eid] = self._count
        self._entity_ids.append(eid)
        self._count += 1

    def search_entities(self, vec: np.ndarray, top_k: int) -> list[tuple[int, float]]:
        """エンティティ・ベクタDB に対する top_k 近傍検索（Algorithm 1 Line 1）。"""
        if self._count == 0:
            return []
        sims = self._entity_vecs[: self._count] @ _normalize(vec).astype(np.float32)
        order = np.argsort(-sims)[:top_k]
        return [(self._entity_ids[i], float(sims[i])) for i in order]

    def find_entity_by_name(self, name: str) -> int | None:
        name_l = name.strip().lower()
        for e in self.entities.values():
            if e.name.strip().lower() == name_l:
                return e.id
        return None

    # ---- 永続化 ----
    def persist(self, storage_dir: str | Path) -> None:
        d = Path(storage_dir)
        d.mkdir(parents=True, exist_ok=True)
        payload = {
            "nodes": [asdict(n) for n in self.nodes.values()],
            "roots": self.roots,
            "entities": [asdict(e) for e in self.entities.values()],
            "relations": self.relations,
            "node_seq": self._node_seq,
            "ent_seq": self._ent_seq,
            "entity_ids": self._entity_ids,
        }
        # アトミック書き込み: 一時ファイルへ書いてから os.replace（中断時の索引破損防止）
        import os

        tmp_json = d / "bookindex.json.tmp"
        tmp_json.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
        os.replace(tmp_json, d / "bookindex.json")
        tmp_npy = d / "entity_vecs.tmp.npy"
        np.save(tmp_npy, self._entity_vecs[: self._count])
        os.replace(tmp_npy, d / "entity_vecs.npy")

    @classmethod
    def load(cls, storage_dir: str | Path) -> "BookIndex":
        d = Path(storage_dir)
        payload = json.loads((d / "bookindex.json").read_text(encoding="utf-8"))
        bi = cls()
        bi.nodes = {n["id"]: TreeNode(**n) for n in payload["nodes"]}
        bi.roots = payload["roots"]
        bi.entities = {e["id"]: Entity(**e) for e in payload["entities"]}
        bi.relations = [tuple(r) for r in payload["relations"]]
        bi._node_seq = payload["node_seq"]
        bi._ent_seq = payload["ent_seq"]
        bi._entity_ids = list(payload["entity_ids"])
        vecs = np.load(d / "entity_vecs.npy")
        bi._entity_vecs = vecs.astype(np.float32)
        bi._count = vecs.shape[0]
        bi._dim = vecs.shape[1] if vecs.ndim == 2 and vecs.shape[0] else None
        bi._row_of = {eid: i for i, eid in enumerate(bi._entity_ids)}
        return bi


def _normalize(v: np.ndarray) -> np.ndarray:
    v = np.asarray(v, dtype=np.float32)
    n = np.linalg.norm(v)
    return v / n if n else v


def render_tree(bi: "BookIndex", *, max_chars: int = 42, show_entities: bool = False) -> str:
    """BookIndex の木を ASCII アートで可視化する（構築結果の確認用）。

    例::
        Handbook
        ├─ [S1] 第1章 概要 (p.1)
        │   ├─ (Text#3 p.1, 812字) 本章では…
        │   └─ (Tbl#4 p.2) モデル | 精度 | 速度
        └─ [S1] 第2章 手順
    """
    rev = bi.node_to_entities() if show_entities else {}

    def _label(n: TreeNode) -> str:
        page = f" p.{n.page}" if n.page is not None else ""
        ents = f" ◆{len(rev[n.id])}エンティティ" if show_entities and n.id in rev else ""
        if n.type == "Section":
            lvl = f"S{n.level}" if n.level else "S"
            return f"[{lvl}] {n.title or n.content}{page}{ents}"
        head = (n.content or "").replace("\n", " ")[:max_chars]
        kind = {"Text": "Text", "Table": "Tbl", "Image": "Img"}.get(n.type, n.type)
        size = f", {len(n.content)}字" if n.type == "Text" else ""
        return f"({kind}#{n.id}{page}{size}) {head}"

    lines: list[str] = []

    def _walk(nid: int, prefix: str, is_last: bool, is_root: bool) -> None:
        n = bi.nodes[nid]
        if is_root:
            lines.append(n.title or n.content)
        else:
            lines.append(prefix + ("└─ " if is_last else "├─ ") + _label(n))
        child_prefix = "" if is_root else prefix + ("    " if is_last else "│   ")
        for i, c in enumerate(n.children):
            _walk(c, child_prefix, i == len(n.children) - 1, False)

    for r in bi.roots:
        _walk(r, "", True, True)
    return "\n".join(lines) if lines else "(空のインデックスです。add_book を実行してください)"


# --------------------------------------------------------------------------
# 4.2 Tree Construction
# --------------------------------------------------------------------------

_MD_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
# 章番号風の見出し: "1", "1.2", "第3章" など
_NUM_HEADING = re.compile(r"^(\d+(\.\d+){0,3})\s+\S|^第[0-9一二三四五六七八九十百]+[章節編]")


def parse_blocks(path: Path, *, ocr=False, layout=False, vlm=False,
                 vlm_model: str | None = None) -> list[dict]:
    """Layout Parsing（4.2.1）。文書を素朴なブロック列に分解する。

    各ブロック: {"content", "type"(Title/Text/Table/Image), "page", "font"}
    - .md/.markdown: 見出しレベルが明示的なので最も正確に木を作れる（推奨）。
    - .pdf: pypdf でページ毎テキストを取り、見出しらしさをヒューリスティック判定。
    - .docx: 見出しスタイル(Heading n)から階層を復元（比較的良好）。
    - .pptx / .xlsx: 一応読めるが BookRAG と相性が悪い（警告あり）。
    - .txt 等: 段落を Text ブロックに。
    """
    suffix = path.suffix.lower()
    if suffix in (".md", ".markdown"):
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            log("UTF-8 で読めないため文字化けを置換して読み込みます（CP932 等は UTF-8 保存を推奨）")
            text = path.read_text(encoding="utf-8", errors="replace")
        return _parse_markdown(text)
    if suffix == ".pdf":
        # 版面解析つきパース（pymupdf のフォントサイズ見出し + 任意 OCR + 任意 VLM 図理解）。
        # pymupdf 未導入 or 全オプション無効なら pypdf のヒューリスティックへ。
        if ocr or vlm or layout not in (None, False):
            from . import docparse

            blocks = docparse.parse_pdf(path, ocr=ocr, layout=(layout or "auto"),
                                        vlm=vlm, vlm_model=vlm_model)
            if blocks is not None:
                if not blocks:
                    log("PDF からテキストを取得できませんでした（スキャンPDFなら ocr=True と "
                        "Tesseract 導入が必要）。")
                return blocks
            log("pymupdf 未導入のため pypdf で解析します（OCR/版面解析/図理解は無効）。"
                "  pip install -e \".[ocr]\"")
        return _parse_pdf(path)
    if suffix in (".docx", ".doc"):
        _warn_format(suffix)
        return _parse_docx(path)
    if suffix == ".pptx":
        _warn_format(suffix)
        return _parse_pptx(path)
    if suffix in (".xlsx", ".xls"):
        _warn_format(suffix)
        return _parse_xlsx(path)
    return _parse_plain(path.read_text(encoding="utf-8", errors="ignore"))


def _warn_format(suffix: str) -> None:
    """BookRAG と相性が悪い形式の注意（デメリットと推奨方式）を表示する。"""
    if suffix in (".docx", ".doc"):
        log("⚠ Word を BookRAG で処理します。見出しスタイル(Heading n)から階層を作りますが、"
            "スタイル未使用の文書は木が浅くなります。構造重視なら Markdown 化、"
            "単純な検索用途なら PagedRAG も検討してください。")
    elif suffix == ".pptx":
        log("⚠ PowerPoint を BookRAG で処理します。スライドは階層的な散文ではないため"
            "木/知識グラフの精度が落ちがちです。推奨: 検索は PagedRAG、要点抽出はチャット。")
    elif suffix in (".xlsx", ".xls"):
        log("⚠ Excel を BookRAG で処理します。表データは BookRAG の木/知識グラフと相性が悪く、"
            "精度が低い可能性が高いです。推奨: 集計・計算は TableQA、検索は PagedRAG、"
            "文章と表が混在するなら DocQA。")


def _parse_docx(path: Path) -> list[dict]:
    """Word。見出しスタイルを Title(level)、本文を Text、表を Table に。"""
    try:
        from docx import Document
    except ImportError:
        log(".docx には python-docx が必要です: pip install python-docx（スキップ）")
        return []
    blocks: list[dict] = []
    try:
        doc = Document(str(path))
    except Exception as e:  # noqa: BLE001
        log(f"Word の読み込みに失敗（スキップ）: {e}")
        return []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.lower().startswith("heading"):
            m = re.search(r"(\d+)", style)
            level = int(m.group(1)) if m else 1
            blocks.append({"content": text, "type": "Title", "page": None, "font": None, "level": level})
        else:
            blocks.append({"content": text, "type": "Text", "page": None, "font": None, "level": None})
    for tb in doc.tables:
        rows = [" | ".join(c.text for c in row.cells) for row in tb.rows]
        rows = [r for r in rows if r.strip()]
        if rows:
            blocks.append({"content": "\n".join(rows), "type": "Table",
                           "page": None, "font": None, "level": None})
    return blocks


def _parse_pptx(path: Path) -> list[dict]:
    """PowerPoint。各スライドのタイトルを Title(level1)、本文を Text に。"""
    try:
        from pptx import Presentation
    except ImportError:
        log(".pptx には python-pptx が必要です: pip install python-pptx（スキップ）")
        return []
    blocks: list[dict] = []
    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        log(f"PowerPoint の読み込みに失敗（スキップ）: {e}")
        return []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:  # noqa: BLE001
            title = None
        blocks.append({"content": title or f"Slide {i}", "type": "Title",
                       "page": i, "font": None, "level": 1})
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = shape.text_frame.text.strip()
            if txt and txt != title:
                blocks.append({"content": txt, "type": "Text", "page": i, "font": None, "level": None})
    return blocks


def _parse_xlsx(path: Path) -> list[dict]:
    """Excel。各シートを Title(level1)、行をまとめて Table に（相性は悪い）。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        log(".xlsx には openpyxl が必要です: pip install openpyxl（スキップ）")
        return []
    blocks: list[dict] = []
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:  # noqa: BLE001
        log(f"Excel の読み込みに失敗（スキップ）: {e}")
        return []
    for ws in wb.worksheets:
        blocks.append({"content": str(ws.title), "type": "Title",
                       "page": str(ws.title), "font": None, "level": 1})
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append({"content": "\n".join(rows), "type": "Table",
                           "page": str(ws.title), "font": None, "level": None})
    return blocks


_MD_TABLE_ROW = re.compile(r"^\s*\|.*\|\s*$")


def _parse_markdown(text: str) -> list[dict]:
    blocks: list[dict] = []
    buf: list[str] = []

    def flush():
        # 本文と表行を分離して emit する。節全体を1つの Table に併合しない
        # （併合すると本文がラベル誤りになり、チャンク化も迂回される）。
        if not buf:
            return
        runs: list[tuple[bool, list[str]]] = []  # (is_table, lines)
        for line in buf:
            is_tbl = bool(_MD_TABLE_ROW.match(line))
            if runs and runs[-1][0] == is_tbl:
                runs[-1][1].append(line)
            else:
                runs.append((is_tbl, [line]))
        for is_tbl, lines in runs:
            content = "\n".join(lines).strip()
            if content:
                blocks.append({"content": content, "type": "Table" if is_tbl else "Text",
                               "page": None, "font": None, "level": None})
        buf.clear()

    for line in text.splitlines():
        m = _MD_HEADING.match(line)
        if m:
            flush()
            blocks.append({"content": m.group(2).strip(), "type": "Title",
                           "page": None, "font": float(7 - len(m.group(1))), "level": len(m.group(1))})
        else:
            buf.append(line)
    flush()
    return blocks


def _parse_plain(text: str) -> list[dict]:
    blocks = []
    for para in re.split(r"\n\s*\n", text):
        para = para.strip()
        if not para:
            continue
        is_heading = bool(_NUM_HEADING.match(para)) and len(para) < 80
        blocks.append({"content": para, "type": "Title" if is_heading else "Text",
                       "page": None, "font": None, "level": None})
    return blocks


def _parse_pdf(path: Path) -> list[dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    blocks = []
    for pno, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for para in re.split(r"\n\s*\n", text):
            para = " ".join(para.split())
            if not para:
                continue
            short = len(para) < 80
            is_heading = short and (bool(_NUM_HEADING.match(para)) or para.isupper())
            blocks.append({"content": para, "type": "Title" if is_heading else "Text",
                           "page": pno, "font": None, "level": None})
    return blocks


def section_filter(blocks: list[dict], *, use_llm: bool = True) -> list[dict]:
    """Section Filtering（4.2.2）。Title 候補に階層 level と最終 type を割り当てる。

    Markdown 由来で level が既知ならそのまま使う。未知（PDF/plain）の Title 候補は
    LLM でレベル判定・誤検出の Text 再分類を行う（use_llm=True）。
    """
    title_idx = [i for i, b in enumerate(blocks) if b["type"] == "Title" and b.get("level") is None]
    candidate_set = set(title_idx)  # LLM が返す id は候補集合のみ受理（無関係ブロックの破壊防止）
    if title_idx and use_llm:
        candidates = [{"id": i, "text": blocks[i]["content"][:120]} for i in title_idx]
        # 長文書対策でバッチ分割。論文 4.2.2 に従い、確定済みの上位レベル節を
        # コンテキストとして次バッチへ持ち回る（バッチ間で階層判定の一貫性を保つ）。
        outline: list[str] = []
        for start in range(0, len(candidates), 40):
            batch = candidates[start : start + 40]
            ctx = ""
            if outline:
                ctx = ("\n\nDocument outline confirmed so far (higher-level sections, in order):\n"
                       + "\n".join(outline[-12:]))
            result = llm_json(
                _PROMPT_SECTION_FILTER + ctx
                + "\n\nCandidates JSON:\n" + json.dumps(batch, ensure_ascii=False)
            )
            if isinstance(result, list):
                for item in result:
                    if not isinstance(item, dict):  # LLM 応答の形状不正はスキップ
                        continue
                    i = item.get("id")
                    if not isinstance(i, int) or i not in candidate_set:
                        continue
                    lvl = item.get("level")
                    if lvl in (None, 0, "None", "none"):
                        blocks[i]["type"] = "Text"  # 誤検出を Text に再分類
                    else:
                        try:
                            level = max(1, int(lvl))
                            blocks[i]["level"] = level
                            if level <= 2:  # 上位レベル節をアウトラインへ
                                outline.append(f"{'#' * level} {blocks[i]['content'][:60]}")
                        except (TypeError, ValueError):
                            pass  # フォールバック（下の番号推定）に任せる
    # LLM 未使用 or 失敗時のフォールバック: 番号の深さからレベル推定
    for i in title_idx:
        b = blocks[i]
        if b["type"] == "Title" and b.get("level") is None:
            m = re.match(r"^(\d+(\.\d+){0,3})", b["content"])
            b["level"] = (m.group(1).count(".") + 1) if m else 1
    return blocks


def blocks_range(blocks):  # 小ヘルパ（id 範囲チェック）
    return range(len(blocks))


def _chunk_text(text: str, chunk_chars: int) -> list[str]:
    """テキストを chunk_chars 程度の塊に分割（改行境界優先）。

    PDF 抽出では段落全体が改行なしの1行になることがあるため、
    1行が上限を超える場合は文字数で強制分割する（従来は無制限に膨らんでいた）。
    """
    text = text.strip()
    if len(text) <= chunk_chars:
        return [text] if text else []
    lines: list[str] = []
    for raw in text.split("\n"):
        while len(raw) > chunk_chars:
            lines.append(raw[:chunk_chars])
            raw = raw[chunk_chars:]
        lines.append(raw)
    chunks, buf, size = [], [], 0
    for line in lines:
        if size + len(line) > chunk_chars and buf:
            chunks.append("\n".join(buf).strip())
            buf, size = [], 0
        buf.append(line)
        size += len(line) + 1
    if buf:
        chunks.append("\n".join(buf).strip())
    return [c for c in chunks if c]


def build_tree(bi: BookIndex, blocks: list[dict], book: str, *, chunk_chars: int = 1500) -> int:
    """フィルタ済みブロック列から木 T を組み立て、その書のルート node id を返す（4.2.2 末尾）。

    本文（Text）は段落ごとにノード化せず、**節の下でまとめて chunk_chars 程度のチャンク**に
    する。これによりノード数（= 後段の LLM 抽出回数）を大幅に削減する。
    Table / Image は構造を保つため単独ノードのまま。
    """
    root = bi.add_node(type="Section", content=book, book=book, level=0, title=book)
    bi.roots.append(root.id)
    stack: list[tuple[int, int]] = [(0, root.id)]

    buf: list[str] = []
    buf_page = [None]  # チャンク先頭ページ

    def _flush(parent: int) -> None:
        if not buf:
            return
        for chunk in _chunk_text("\n".join(buf), chunk_chars):
            node = bi.add_node(type="Text", content=chunk, book=book,
                               page=buf_page[0], parent=parent)
            bi.nodes[parent].children.append(node.id)
        buf.clear()
        buf_page[0] = None

    for b in blocks:
        parent = stack[-1][1] if stack else root.id
        if b["type"] == "Title":
            _flush(parent)
            level = b.get("level") or 1
            while stack and stack[-1][0] >= level:
                stack.pop()
            parent = stack[-1][1] if stack else root.id
            node = bi.add_node(type="Section", content=b["content"], book=book,
                               level=level, title=b["content"], page=b.get("page"), parent=parent)
            bi.nodes[parent].children.append(node.id)
            stack.append((level, node.id))
        elif b["type"] in ("Table", "Image"):
            _flush(parent)
            node = bi.add_node(type=b["type"], content=b["content"], book=book,
                               page=b.get("page"), parent=parent)
            bi.nodes[parent].children.append(node.id)
        else:  # Text: バッファに溜めてチャンク化
            if buf_page[0] is None:
                buf_page[0] = b.get("page")
            elif b.get("page") != buf_page[0]:
                # ページが変わったら区切る。見出しの無い文書で全チャンクが
                # 先頭ページ番号になる（ページ出典の喪失）のを防ぐ。
                _flush(parent)
                buf_page[0] = b.get("page")
            buf.append(b["content"])

    _flush(stack[-1][1] if stack else root.id)
    return root.id


# --------------------------------------------------------------------------
# 4.3 Graph Construction + Gradient-based Entity Resolution (Algorithm 1)
# --------------------------------------------------------------------------


def _section_ancestor_id(bi: "BookIndex", nid: int) -> int:
    """ノードが属する最も近い Section 祖先（無ければ自分）の id。"""
    cur = nid
    seen = 0
    while bi.nodes[cur].type != "Section" and bi.nodes[cur].parent is not None and seen < 10000:
        cur = bi.nodes[cur].parent
        seen += 1
    return cur


def _even_sample(bi: "BookIndex", targets: list[int], budget: int) -> list[int]:
    """max_nodes を超える場合、セクションごとに均等（round-robin）にサンプリングする。

    従来は先頭から budget 件だけ採用していたため、後半セクションが丸ごと無視され
    「前半しか見ない」偏りが出ていた。文書全体を薄く広くカバーするよう変更。
    後段の Gradient ER は順序依存のため、**返り値は文書順（targets の元の並び）**を保つ。
    """
    from collections import OrderedDict
    from itertools import zip_longest

    if budget >= len(targets):
        return list(targets)
    buckets: "OrderedDict[int, list[int]]" = OrderedDict()
    for nid in targets:  # 入力順（=文書順）を保ったままセクション別に束ねる
        buckets.setdefault(_section_ancestor_id(bi, nid), []).append(nid)
    # ラウンドロビンで各セクションから均等に採用（O(total)）
    picked: set[int] = set()
    for row in zip_longest(*buckets.values()):
        for nid in row:
            if nid is not None and len(picked) < budget:
                picked.add(nid)
        if len(picked) >= budget:
            break
    return [nid for nid in targets if nid in picked]  # 文書順を復元


def build_graph(bi: BookIndex, node_ids: list[int], *, gradient_g: float = 0.6,
                er_top_k: int = 10, max_workers: int = 8, min_chars: int = 40,
                max_nodes: int = 300, er_use_llm: bool = False, reranker=None,
                all_nodes: bool = False) -> None:
    """各ノードからエンティティ/関係を抽出し、Gradient-based ER で KG を構築する。

    速度のため、抽出（LLM 呼び出し + 埋め込み）はノード単位で **並列化** し、
    Gradient ER とグラフ構築（順序依存・BookIndex を変更する）は **逐次** で行う。
    - 短すぎるノード（min_chars 未満）は抽出をスキップ。
    - all_nodes=True なら全ノードを処理。False（既定）で max_nodes を超える場合は
      **セクションごとに均等サンプリング**（先頭打ち切りにせず文書全体をカバー）。
    - er_use_llm=False（既定）なら曖昧マージで LLM を呼ばず、最有力候補を採用（高速）。
    - reranker 指定時は Algorithm 1 の Rerank model R として名寄せ候補を再スコアする。
    - Table ノードは論文 4.3.1 に従い v_table エンティティ + ヘッダを ContainedIn で構造化。
    """
    from concurrent.futures import ThreadPoolExecutor

    targets = [
        nid for nid in node_ids
        if bi.nodes[nid].type != "Section"
        and len(bi.nodes[nid].content.strip()) >= min_chars
    ]
    if not targets:
        # 空文書/スキャンPDF等。従来はここで黙って戻り「読めていない」ことに気づけなかった。
        log("警告: 抽出対象の本文ノードが 0 件です。文書からテキストを取得できていない可能性が"
            "あります（スキャンPDF・空文書・全ノードが min_chars 未満）。")
        return
    if not all_nodes and len(targets) > max_nodes:
        print(f"[BookRAG] 対象ノード {len(targets)} 件を max_nodes={max_nodes} 件へ"
              "**セクション均等サンプリング**で圧縮します（全件は all_nodes=True）")
        targets = _even_sample(bi, targets, max_nodes)

    def _work(nid: int, fail_fast: bool = True):
        # 抽出フェーズは BookIndex を読むだけ（スレッド安全）。失敗しても全体を止めず、
        # 状態（ok / empty / badjson / error:*）を返して後段で集計・再試行する。
        try:
            node = bi.nodes[nid]
            data = _extract_graph(node, fail_fast=fail_fast)   # entities + relations を1回で
            if data is None:
                return nid, [], [], None, "badjson"            # モデルが JSON を返さない
            ents = data["entities"]
            if not ents:
                return nid, [], [], None, "empty"
            names = [f"{e['name']} ({e.get('type','')}): {e.get('description','')}" for e in ents]
            return nid, ents, data["relations"], embed(names), "ok"
        except Exception as e:  # noqa: BLE001
            return nid, [], [], None, f"error:{type(e).__name__}"

    results = []
    total = len(targets)
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        for res in progress(ex.map(_work, targets), total=total,
                            desc=f"抽出: エンティティ/関係 ({max_workers}並列)"):
            results.append(res)

    # 失敗ノード（通信エラー/JSON不正）を逐次(1並列)で1回だけ再試行する。
    # ローカルLLMサーバが並列リクエストを捌けずタイムアウトする構成の救済。
    retry_idx = [i for i, r in enumerate(results)
                 if r[4].startswith("error") or r[4] == "badjson"]
    if retry_idx:
        log(f"{len(retry_idx)}/{len(results)} ノードの抽出が失敗。サーバ過負荷の可能性があるため"
            "逐次(1並列)で再試行します…")
        for i in progress(retry_idx, total=len(retry_idx), desc="失敗ノードの再試行 (1並列)"):
            results[i] = _work(results[i][0], fail_fast=False)

    # 状態を集計して報告（従来は失敗が黙殺され「エンティティ0」の原因が見えなかった）。
    from collections import Counter

    stat = Counter(("error" if r[4].startswith("error") else r[4]) for r in results)
    errors = Counter(r[4].split(":", 1)[1] for r in results if r[4].startswith("error"))
    found = sum(len(r[1]) for r in results)
    log(f"抽出結果: 成功 {stat.get('ok', 0)} / エンティティ0 {stat.get('empty', 0)} / "
        f"JSON不正 {stat.get('badjson', 0)} / 通信エラー {stat.get('error', 0)}"
        f"（エンティティ合計 {found}）")
    if stat.get("error"):
        log(f"通信エラー内訳: {dict(errors)}。ローカルLLMサーバが並列を捌けない場合は "
            "BookRAG(max_workers=1) を、1リクエストが遅い場合は "
            "configure(request_timeout=300) を試してください。")
    if stat.get("badjson"):
        log("モデルが JSON 形式で応答していません。"
            "llmlab.complete('Return JSON: {\"ok\":true}') で挙動を確認してください。"
            "思考過程を出力するモデル（Qwen3/R1系）の場合、思考が長すぎて JSON が"
            "途中で切れている可能性があります。サーバ側で思考の無効化"
            "（enable_thinking=false や /no_think）を推奨します（速度も改善）。")
    if stat.get("empty", 0) >= max(1, len(results) // 2):
        log("半数以上のノードで抽出結果が空でした。思考過程を出力するモデルは思考で"
            "トークン上限に達し空の JSON を返しがちです。サーバ側で思考の無効化を"
            "推奨します。改善しない場合は指示追従の強いモデルの利用を検討してください。")
    if found == 0:
        log("警告: エンティティが1件も抽出できませんでした（上記の内訳を参照）。")

    # Gradient ER + グラフ構築は順序依存のため逐次（入力ノード順）。進捗も表示。
    for nid, ents, rels, vecs, _st in progress(results, total=len(results),
                                               desc="名寄せ/グラフ構築 (KG)"):
        node = bi.nodes[nid]
        local_name_to_eid: dict[str, int] = {}
        if ents:
            for e, vec in zip(ents, vecs):
                eid = gradient_entity_resolution(
                    bi, e["name"], e.get("type", ""), e.get("description", ""),
                    nid, vec, gradient_g=gradient_g, er_top_k=er_top_k, use_llm=er_use_llm,
                    reranker=reranker,
                )
                local_name_to_eid[e["name"].strip().lower()] = eid
            for rel in rels:
                s = local_name_to_eid.get(str(rel.get("source", "")).strip().lower())
                t = local_name_to_eid.get(str(rel.get("target", "")).strip().lower())
                if s is not None and t is not None and s != t:
                    bi.relations.append((s, t, str(rel.get("relation", "related_to"))))
        # 論文 4.3.1: Table ノードは v_table エンティティ + ヘッダ(ContainedIn) を構造化。
        # LLM 抽出が空でも実行する（プロンプト任せにせず構造として保証する）。
        if node.type == "Table":
            try:
                _augment_table_entities(
                    bi, node, local_name_to_eid,
                    gradient_g=gradient_g, er_top_k=er_top_k,
                    use_llm=er_use_llm, reranker=reranker,
                )
            except Exception as e:  # noqa: BLE001
                log(f"表エンティティの構造化に失敗（スキップ）: {e}")


def _augment_table_entities(bi: BookIndex, node: TreeNode, local_name_to_eid: dict,
                            *, gradient_g: float, er_top_k: int, use_llm: bool,
                            reranker=None) -> None:
    """論文 4.3.1 の Table 処理: v_table を作り、行/列ヘッダと抽出エンティティを
    ContainedIn で v_table に接続する。"""
    rows = [r for r in node.content.split("\n") if r.strip()]
    if not rows:
        return
    header_cells = [c.strip() for c in rows[0].split("|") if c.strip()][:10]
    first_col = [r.split("|")[0].strip() for r in rows[1:16] if r.split("|")[0].strip()]
    page = f" p.{node.page}" if node.page is not None else ""
    table_name = f"表({node.book}{page}): " + (" / ".join(header_cells[:4]) or rows[0][:40])

    headers = list(dict.fromkeys(header_cells + first_col))  # 重複排除・順序維持
    names = [table_name] + headers
    vecs = embed(names)
    v_table = gradient_entity_resolution(
        bi, table_name, "Table", node.content[:160], node.id, vecs[0],
        gradient_g=gradient_g, er_top_k=er_top_k, use_llm=use_llm, reranker=reranker,
    )
    for nm, vec in zip(headers, vecs[1:]):
        eid = gradient_entity_resolution(
            bi, nm, "TableHeader", f"{table_name} のヘッダ", node.id, vec,
            gradient_g=gradient_g, er_top_k=er_top_k, use_llm=use_llm, reranker=reranker,
        )
        if eid != v_table:
            bi.relations.append((eid, v_table, "ContainedIn"))
    # そのノードから LLM 抽出されたエンティティも v_table に接続（論文: primary vertex に link）
    for eid in set(local_name_to_eid.values()):
        if eid != v_table:
            bi.relations.append((eid, v_table, "ContainedIn"))


def gradient_entity_resolution(bi: BookIndex, name: str, etype: str, desc: str, origin: int,
                               vec: np.ndarray, *, gradient_g: float, er_top_k: int,
                               use_llm: bool = False, reranker=None) -> int:
    """Algorithm 1（Gradient-based entity resolution）。返り値は確定したエンティティ id。

    新規概念なら追加、既存の別名なら最も確からしい正準エンティティへマージする。
    reranker 指定時はそれを Rerank model R として候補を再スコア（論文どおり）。
    未指定時はベクトルDB のコサイン類似度を rerank スコア S とみなす。
    use_llm=True のときだけ、複数の高関連候補がある場合に LLM で1つ選ぶ（遅い）。
    """
    candidates = bi.search_entities(vec, er_top_k)  # Line 1: Search
    if not candidates:
        return bi.add_entity(name, etype, desc, origin, vec).id

    # ほぼ同一（コサイン>=0.95）は勾配判定を待たず即マージする。
    # 候補が er_top_k 未満の小規模索引では「リストを使い切った」ことが Case A と
    # 区別できず、完全一致でも重複追加されてしまうため（索引成長初期の名寄せ崩壊対策）。
    if candidates[0][1] >= 0.95:
        bi.merge_entity(candidates[0][0], name, origin, vec)
        return candidates[0][0]

    # Line 2-3: Rerank model R による再スコア（reranker 指定時）→ 降順ソート
    candidates = _rerank_candidates(bi, name, etype, desc, candidates, reranker)
    scores = [s for _, s in candidates]
    sel = [candidates[0][0]]          # Line 4: Sel ← Ec[0]
    prev = scores[0]
    for (eid, sc) in candidates[1:]:  # Line 5-8: 連続スコアの急落（gradient）を検出
        if sc > gradient_g * prev:    # 緩やかな低下 → 高関連集合に含める
            sel.append(eid)
            prev = sc
        else:
            break

    if len(sel) == len(candidates):   # Line 9-10: Case A（勾配なし）→ 新規エンティティ
        return bi.add_entity(name, etype, desc, origin, vec).id

    # Line 11-14: Case B → マージ
    if len(sel) == 1 or not use_llm:
        v_sel = sel[0]                # 既定: LLM を呼ばず最有力候補を採用（高速）
    else:
        v_sel = _llm_select_entity(bi, name, etype, desc, sel)
        if v_sel is None:
            return bi.add_entity(name, etype, desc, origin, vec).id
    bi.merge_entity(v_sel, name, origin, vec)
    return v_sel


def _rerank_candidates(bi: BookIndex, name: str, etype: str, desc: str,
                       candidates: list[tuple[int, float]], reranker):
    """Algorithm 1 Line 2-3: Rerank model R で候補を再スコアし降順に返す。

    reranker 未指定（または cosine）ならベクトル検索スコアのまま。CrossEncoder の
    ロジット等 [0,1] 外のスコアは sigmoid で写像する（勾配判定は比率ベースのため
    正の単調スコアを前提とする）。失敗時はコサインスコアで続行。
    """
    if reranker is None or type(reranker).__name__ == "CosineReranker":
        return candidates
    try:
        query = f"{name} ({etype}): {desc[:160]}"
        docs = [
            f"{bi.entities[eid].name} ({bi.entities[eid].type}): {bi.entities[eid].description[:160]}"
            for eid, _ in candidates
        ]
        scores = reranker.rerank(query, docs)
        if any(s < 0.0 or s > 1.0 for s in scores):
            import math

            scores = [1.0 / (1.0 + math.exp(-s)) for s in scores]
        return sorted(zip((eid for eid, _ in candidates), scores), key=lambda x: -x[1])
    except Exception as e:  # noqa: BLE001
        log(f"ER rerank に失敗（コサインで続行）: {e}")
        return candidates


def _extract_graph(node: TreeNode, *, fail_fast: bool = False) -> dict | None:
    """1 回の LLM 呼び出しでノードからエンティティと関係をまとめて抽出する。

    fail_fast=True で SDK リトライ無効（並列パスで使用。失敗は上位の逐次再試行に回す）。
    モデルが JSON を返さなかった場合は None（「エンティティ0」と区別するため）。
    """
    result = llm_json(
        _PROMPT_GRAPH_EXTRACT + f"\n\nNode type: {node.type}\nContent:\n{node.content[:2500]}",
        max_retries=0 if fail_fast else None,
        # 長時間生成→タイムアウトを防ぐ上限。思考過程を出すモデルは思考にも消費するため
        # 並列パス2000 / 逐次レスキュー4000（思考でトークンを使い切り空 JSON になる対策）
        max_tokens=2000 if fail_fast else 4000,
    )
    if result is None:
        return None  # JSON 解釈不能（指示無視）
    ents: list[dict] = []
    rels: list[dict] = []
    if isinstance(result, dict):
        raw_e = result.get("entities", [])
        if isinstance(raw_e, list):
            ents = [e for e in raw_e if isinstance(e, dict) and e.get("name")]
        raw_r = result.get("relations", [])
        if isinstance(raw_r, list):
            rels = [r for r in raw_r if isinstance(r, dict) and r.get("source") and r.get("target")]
    return {"entities": ents, "relations": rels}


def _llm_select_entity(bi: BookIndex, name: str, etype: str, desc: str,
                       candidate_ids: list[int]) -> int | None:
    """Algorithm 1 Line 13: 複数の高関連候補から正準エンティティを LLM で1つ選ぶ。"""
    cands = [
        {"id": eid, "name": bi.entities[eid].name, "type": bi.entities[eid].type,
         "description": bi.entities[eid].description[:160]}
        for eid in candidate_ids
    ]
    result = llm_json(
        _PROMPT_ENTITY_RESOLUTION
        + f"\n\nNew Entity: {json.dumps({'name': name, 'type': etype, 'description': desc[:160]}, ensure_ascii=False)}"
        + f"\nCandidate Entities: {json.dumps(cands, ensure_ascii=False)}"
    )
    if isinstance(result, dict):
        sid = result.get("select_id", -1)
        try:
            sid = int(sid)
        except (TypeError, ValueError):
            return None
        return sid if sid in candidate_ids else None
    return None


# --------------------------------------------------------------------------
# プロンプト（論文 Fig.10-13 に対応／日本語文書にも対応するよう調整）
# --------------------------------------------------------------------------

_PROMPT_SECTION_FILTER = """You analyze heading candidates extracted from a document and decide,
for each, its hierarchical level (1 = top-level section, 2, 3, ...) or null if it is NOT a real
section heading (i.e., it is body text mistakenly detected as a title).
Return ONLY a JSON array: [{"id": <int>, "level": <int or null>}, ...]."""

_PROMPT_GRAPH_EXTRACT = """You are an information extraction engine. From the given document node,
extract (1) AT MOST 10 key entities and (2) the directed relations among those entities (relations
must only reference the extracted entities). Keep names canonical, descriptions under 15 words each,
and output in the document's language.
Return ONLY JSON (no explanations, no markdown fences):
{"entities":[{"name": str, "type": str, "description": str}, ...],
 "relations":[{"source": str, "target": str, "relation": str}, ...]}."""

_PROMPT_ENTITY_RESOLUTION = """You are an Entity Resolution Adjudicator. Decide if the New Entity refers to
the EXACT same real-world concept as one of the Candidate Entities. Be strict and conservative; when in
doubt output -1. Names must be extremely similar, a direct abbreviation, or a well-known alias. Distinct
parallel concepts are NOT a match even if descriptions look similar.
Return ONLY JSON: {"select_id": <candidate id or -1>, "explanation": str}."""
